"""Extract text from common office and text-format attachments.

Strategy
--------
We don't try to OCR images, render slides, or parse complex layouts. The
goal is to turn predictable formats into plain text the agent can reason
about. PDFs and images go through Claude itself (vision + PDF native), not
through this module — they're returned as `kind="multimodal"` so the
adapter passes the path to the agent and lets Claude's Read tool handle
the rest.

Supported extraction:
- .docx (python-docx)
- .xlsx, .xlsm (openpyxl)
- .pptx (python-pptx)
- .txt, .md, .csv, .json, .yaml, .yml, .py, .js, .ts (read as utf-8)

Returned as `kind="multimodal"` (delegate to Claude Read):
- .pdf, .png, .jpg, .jpeg, .gif, .webp

`kind="unknown"`: everything else. Adapter can still pass the path to the
agent — Claude Read may handle some formats we didn't anticipate.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Literal

logger = logging.getLogger(__name__)

FILE_KIND = Literal["text-extracted", "multimodal", "unknown"]


_TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".py", ".js", ".ts", ".log"}
_OFFICE_EXTS = {".docx", ".xlsx", ".xlsm", ".pptx"}
_MULTIMODAL_EXTS = {".pdf", ".png", ".jpg", ".jpeg", ".gif", ".webp"}


def classify_path(path: Path) -> FILE_KIND:
    suffix = path.suffix.lower()
    if suffix in _TEXT_EXTS or suffix in _OFFICE_EXTS:
        return "text-extracted"
    if suffix in _MULTIMODAL_EXTS:
        return "multimodal"
    return "unknown"


def extract_text(path: Path) -> str | None:
    """Extract plain text from an attachment.

    Returns:
        - The extracted text string for office/text formats.
        - None for multimodal formats (caller should hand the path to Claude).
        - None if extraction fails — caller logs and falls back to path-based.
    """
    suffix = path.suffix.lower()

    if suffix in _TEXT_EXTS:
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except Exception as e:
            logger.warning("could not read text file %s: %s", path, e)
            return None

    if suffix == ".docx":
        try:
            from docx import Document  # type: ignore[import-not-found]
        except ImportError:
            logger.warning("python-docx not installed; cannot extract %s", path.name)
            return None
        try:
            doc = Document(str(path))
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            # Tables
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    paras.append(" | ".join(cells))
            return "\n".join(paras)
        except Exception as e:
            logger.exception("docx extraction failed for %s", path)
            return None

    if suffix in {".xlsx", ".xlsm"}:
        try:
            from openpyxl import load_workbook  # type: ignore[import-not-found]
        except ImportError:
            logger.warning("openpyxl not installed; cannot extract %s", path.name)
            return None
        try:
            wb = load_workbook(str(path), data_only=True, read_only=True)
            chunks: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                chunks.append(f"# Sheet: {sheet_name}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        chunks.append(" | ".join(cells))
                chunks.append("")
            return "\n".join(chunks)
        except Exception as e:
            logger.exception("xlsx extraction failed for %s", path)
            return None

    if suffix == ".pptx":
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
        except ImportError:
            logger.warning("python-pptx not installed; cannot extract %s", path.name)
            return None
        try:
            prs = Presentation(str(path))
            chunks: list[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                chunks.append(f"# Slide {i}")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "") or ""
                    if text.strip():
                        chunks.append(text)
                chunks.append("")
            return "\n".join(chunks)
        except Exception as e:
            logger.exception("pptx extraction failed for %s", path)
            return None

    return None
